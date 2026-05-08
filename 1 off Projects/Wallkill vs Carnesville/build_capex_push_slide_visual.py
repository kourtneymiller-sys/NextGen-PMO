from __future__ import annotations

from pathlib import Path
import re

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

BASE = Path(__file__).parent
ROOT = BASE.parent.parent
TEMPLATE = ROOT / "Presentation Template.pptx"
CAPEX_FILE = BASE / "Wallkill CapEx Push.xlsx"
OUT = BASE / "Wallkill_CapEx_Pushout_Visual_One_Slide.pptx"

WMT_BLUE = RGBColor(0, 83, 226)
WMT_BLUE_2 = RGBColor(15, 96, 240)
WMT_DARK = RGBColor(0, 31, 94)
WMT_YELLOW = RGBColor(255, 194, 32)
GREEN = RGBColor(42, 135, 3)
RED = RGBColor(232, 17, 0)
ORANGE = RGBColor(245, 124, 0)
PURPLE = RGBColor(91, 53, 213)
TEAL = RGBColor(0, 130, 148)
GRAY = RGBColor(88, 96, 111)
LIGHT_GRAY = RGBColor(244, 247, 251)
MID_GRAY = RGBColor(213, 222, 235)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 35, 48)
PALE_BLUE = RGBColor(232, 240, 254)
PALE_YELLOW = RGBColor(255, 248, 216)
PALE_RED = RGBColor(255, 235, 232)
PALE_GREEN = RGBColor(238, 248, 235)

FY_COLS = {"FY26": 4, "FY27": 5, "FY28": 6, "FY29": 7, "FY30": 8, "FY31": 9}
ORDER = ["CONSTRUCTION", "ENERGY", "INITIAL ORDER", "MHE", "TECHNOLOGY", "INTERIM INTEREST"]
CAT_COLORS = {
    "Construction": WMT_BLUE,
    "MHE": PURPLE,
    "Technology": TEAL,
    "Energy": ORANGE,
    "Initial Order": GREEN,
    "Interim Interest": RED,
}


def money(v):
    if v is None or v == "":
        return 0.0
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip().replace("$", "").replace(",", "").replace("(", "-").replace(")", "")
    try:
        return float(s) if s else 0.0
    except ValueError:
        return 0.0


def fmt_m(v: float, signed: bool = False) -> str:
    sign = ""
    if signed and v > 0:
        sign = "+"
    elif v < 0:
        sign = "-"
    return f"{sign}${abs(v)/1_000_000:.1f}M"


def norm_label(v) -> str:
    return re.sub(r"\s+", " ", str(v or "")).strip().upper()


def read_rows():
    wb = openpyxl.load_workbook(CAPEX_FILE, data_only=True)
    ws = wb["CapEx Spend by FY (2)"]
    return list(ws.iter_rows(values_only=True))


def parse_profile(rows, start_marker: str) -> dict[str, dict[str, float]]:
    start = None
    for i, row in enumerate(rows):
        if start_marker.lower() in str(row[1] or "").lower():
            start = i
            break
    if start is None:
        raise RuntimeError(f"Could not find section: {start_marker}")
    data = {}
    for row in rows[start + 2:]:
        label = norm_label(row[1] if len(row) > 1 else "")
        if not label:
            break
        if label in {"TOTAL CAPEX BUDGET", "TOTAL EX INTERIM INTEREST", "LAND"}:
            continue
        if label not in ORDER:
            continue
        if label in data:
            continue
        data[label] = {fy: money(row[idx] if idx < len(row) else 0) for fy, idx in FY_COLS.items()}
    return data


def compute():
    rows = read_rows()
    original = parse_profile(rows, "Wallkill  - Original")
    updated = parse_profile(rows, "Wallkill Updated")
    out = []
    for c in ORDER:
        o = original.get(c, {fy: 0 for fy in FY_COLS})
        u = updated.get(c, {fy: 0 for fy in FY_COLS})
        delta = {fy: u.get(fy, 0) - o.get(fy, 0) for fy in FY_COLS}
        early_reduction = -(delta["FY27"] + delta["FY28"])
        future_increase = delta["FY29"] + delta["FY30"] + delta["FY31"]
        name = c.title().replace("Mhe", "MHE")
        out.append({
            "category": name,
            **{k.lower(): v for k, v in delta.items()},
            "early_reduction": early_reduction,
            "future_increase": future_increase,
            "net": future_increase - early_reduction,
        })
    totals = {
        "fy27": sum(r["fy27"] for r in out),
        "fy28": sum(r["fy28"] for r in out),
        "fy29": sum(r["fy29"] for r in out),
        "fy30": sum(r["fy30"] for r in out),
        "fy31": sum(r["fy31"] for r in out),
        "early_reduction": sum(r["early_reduction"] for r in out),
        "future_increase": sum(r["future_increase"] for r in out),
    }
    totals["interest_reduction"] = totals["early_reduction"] - totals["future_increase"]
    return out, totals


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT, valign=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos Display"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    s.line.width = Pt(0.75)
    return s


def add_line(slide, x1, y1, x2, y2, color, width=2.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_arrow(slide, x, y, w, h, fill, text, subtext="", text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos Display"
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = text_color
    if subtext:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtext
        r2.font.name = "Aptos Display"
        r2.font.bold = True
        r2.font.size = Pt(8.5)
        r2.font.color.rgb = text_color
    return shp


def add_kpi(slide, x, y, w, title, value, sub, fill, accent):
    add_rect(slide, x, y, w, Inches(0.9), fill, accent, True)
    add_rect(slide, x, y, Inches(0.08), Inches(0.9), accent, accent, False)
    add_text(slide, x + Inches(0.18), y + Inches(0.09), w - Inches(0.3), Inches(0.18), title, 8.8, True, GRAY)
    add_text(slide, x + Inches(0.18), y + Inches(0.31), Inches(1.45), Inches(0.34), value, 21, True, accent)
    add_text(slide, x + Inches(1.62), y + Inches(0.37), w - Inches(1.78), Inches(0.22), sub, 8.5, False, BLACK)


def add_category_card(slide, x, y, w, h, row):
    color = CAT_COLORS.get(row["category"], WMT_BLUE)
    add_rect(slide, x, y, w, h, WHITE, MID_GRAY, True)
    add_rect(slide, x, y, Inches(0.08), h, color, color, False)
    add_text(slide, x + Inches(0.16), y + Inches(0.08), w - Inches(0.25), Inches(0.18), row["category"], 8.5, True, BLACK)
    add_text(slide, x + Inches(0.16), y + Inches(0.31), Inches(0.9), Inches(0.24), fmt_m(row["early_reduction"]), 13.5, True, RED)
    add_text(slide, x + Inches(1.02), y + Inches(0.35), Inches(0.55), Inches(0.18), "out", 7.5, True, GRAY)
    add_text(slide, x + Inches(1.55), y + Inches(0.31), Inches(0.95), Inches(0.24), fmt_m(row["future_increase"], True), 13.5, True, WMT_BLUE)
    add_text(slide, x + Inches(2.42), y + Inches(0.35), Inches(0.55), Inches(0.18), "later", 7.5, True, GRAY)


def build():
    rows, totals = compute()
    top = sorted(rows, key=lambda r: abs(r["early_reduction"]), reverse=True)

    prs = Presentation(str(TEMPLATE)) if TEMPLATE.exists() else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # full background
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), RGBColor(248, 250, 253), RGBColor(248, 250, 253), False)

    # Header band
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.78), WMT_BLUE, WMT_BLUE, False)
    add_text(slide, Inches(0.36), Inches(0.13), Inches(8.8), Inches(0.42),
             "Wallkill CapEx Pushout: Reduces Near-Term Cash Draw", 20.5, True, WHITE)
    add_text(slide, Inches(9.1), Inches(0.2), Inches(3.85), Inches(0.25),
             "Permitting timeline alignment", 10.8, True, WMT_YELLOW, PP_ALIGN.RIGHT)

    # KPI row
    add_kpi(slide, Inches(0.35), Inches(0.98), Inches(3.9), "FY27/FY28 CapEx pulled down", fmt_m(totals["early_reduction"]), "removed from near-term profile", PALE_RED, RED)
    add_kpi(slide, Inches(4.72), Inches(0.98), Inches(3.9), "Rephased to FY29-FY31", fmt_m(totals["future_increase"]), "matched to extended permit path", PALE_BLUE, WMT_BLUE)
    add_kpi(slide, Inches(9.08), Inches(0.98), Inches(3.9), "Interim interest benefit", fmt_m(totals["interest_reduction"]), "reduced carrying cost", PALE_GREEN, GREEN)

    # Main flow graphic
    add_text(slide, Inches(0.45), Inches(2.08), Inches(12.45), Inches(0.26),
             "CapEx timing shift: from near-term funding pressure to later-year execution window", 13, True, BLACK)

    # Left and right pools
    add_rect(slide, Inches(0.55), Inches(2.58), Inches(2.55), Inches(1.05), PALE_RED, RED, True)
    add_text(slide, Inches(0.72), Inches(2.72), Inches(2.2), Inches(0.22), "FY27/FY28", 12, True, RED, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.72), Inches(2.98), Inches(2.2), Inches(0.34), f"-{fmt_m(totals['early_reduction'])}", 25, True, RED, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.72), Inches(3.34), Inches(2.2), Inches(0.18), "near-term spend removed", 8.3, True, GRAY, PP_ALIGN.CENTER)

    add_arrow(slide, Inches(3.36), Inches(2.66), Inches(2.2), Inches(0.87), WMT_YELLOW, "REPHASE", "protects interim interest")

    # Split later years
    year_vals = [("FY29", totals["fy29"], WMT_BLUE), ("FY30", totals["fy30"], PURPLE), ("FY31", totals["fy31"], TEAL)]
    x = Inches(5.95)
    max_val = max(abs(v) for _, v, _ in year_vals) or 1
    for yr, val, color in year_vals:
        h = Inches(0.45 + 0.58 * (abs(val) / max_val))
        y = Inches(3.62) - h
        add_rect(slide, x, y, Inches(1.35), h, color, color, True)
        add_text(slide, x + Inches(0.06), y + Inches(0.08), Inches(1.23), Inches(0.18), yr, 9.5, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.06), y + Inches(0.31), Inches(1.23), Inches(0.22), fmt_m(val, True), 11.5, True, WHITE, PP_ALIGN.CENTER)
        x += Inches(1.53)
    add_text(slide, Inches(5.95), Inches(3.78), Inches(4.5), Inches(0.2),
             f"Net later-year CapEx increase: {fmt_m(totals['future_increase'])}", 9.2, True, WMT_BLUE, PP_ALIGN.CENTER)

    # Small annual delta chart
    chart_data = CategoryChartData()
    chart_data.categories = ["FY27", "FY28", "FY29", "FY30", "FY31"]
    vals = [totals["fy27"]/1e6, totals["fy28"]/1e6, totals["fy29"]/1e6, totals["fy30"]/1e6, totals["fy31"]/1e6]
    chart_data.add_series("$M", vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(10.55), Inches(2.48), Inches(2.35), Inches(1.6), chart_data).chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(6.5)
    chart.category_axis.tick_labels.font.size = Pt(6.5)
    chart.value_axis.has_major_gridlines = False
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RED if vals[i] < 0 else WMT_BLUE

    # Category cards
    add_text(slide, Inches(0.45), Inches(4.28), Inches(12.45), Inches(0.24),
             "Category-level pushout", 13, True, BLACK)
    y1 = Inches(4.66)
    y2 = Inches(5.55)
    card_w = Inches(4.0)
    for i, row in enumerate(top[:3]):
        add_category_card(slide, Inches(0.45) + Inches(4.18) * i, y1, card_w, Inches(0.68), row)
    for i, row in enumerate(top[3:6]):
        add_category_card(slide, Inches(0.45) + Inches(4.18) * i, y2, card_w, Inches(0.68), row)

    # Bottom takeaway strip
    add_rect(slide, Inches(0.35), Inches(6.55), Inches(12.63), Inches(0.52), WMT_DARK, WMT_DARK, True)
    add_text(slide, Inches(0.58), Inches(6.66), Inches(12.17), Inches(0.26),
             "Takeaway: The updated Wallkill profile pulls $593M out of FY27/FY28, shifts $575M into FY29-FY31, and lowers interim interest by ~$18M while aligning spend to the extended permitting timeline.",
             10.5, True, WHITE, PP_ALIGN.CENTER)

    add_text(slide, Inches(0.4), Inches(7.18), Inches(12.55), Inches(0.16),
             "Source: Wallkill CapEx Push.xlsx; comparison reflects Wallkill Original vs. Wallkill Updated CapEx profile.", 7.0, False, GRAY)

    prs.save(str(OUT))
    print(f"Created: {OUT}")
    print(f"FY27/FY28 reduction: {fmt_m(totals['early_reduction'])}")
    print(f"FY29-FY31 increase: {fmt_m(totals['future_increase'])}")
    print(f"Interim interest benefit: {fmt_m(totals['interest_reduction'])}")


if __name__ == "__main__":
    build()
