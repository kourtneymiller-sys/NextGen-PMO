from __future__ import annotations

from pathlib import Path
import re

import openpyxl
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).parent
ROOT = BASE.parent.parent
TEMPLATE = ROOT / "Presentation Template.pptx"
CAPEX_FILE = BASE / "Wallkill CapEx Push.xlsx"
OUT = BASE / "Wallkill_CapEx_Pushout_One_Slide.pptx"

WMT_BLUE = RGBColor(0, 83, 226)
WMT_DARK = RGBColor(0, 31, 94)
WMT_YELLOW = RGBColor(255, 194, 32)
GREEN = RGBColor(42, 135, 3)
RED = RGBColor(232, 17, 0)
GRAY = RGBColor(88, 96, 111)
LIGHT_GRAY = RGBColor(243, 246, 250)
MID_GRAY = RGBColor(213, 222, 235)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 35, 48)
PALE_BLUE = RGBColor(232, 240, 254)
PALE_YELLOW = RGBColor(255, 248, 216)

FY_COLS = {"FY26": 4, "FY27": 5, "FY28": 6, "FY29": 7, "FY30": 8, "FY31": 9}
EARLY = ["FY27", "FY28"]
FUTURE = ["FY29", "FY30", "FY31"]
ORDER = ["CONSTRUCTION", "ENERGY", "INITIAL ORDER", "MHE", "TECHNOLOGY", "INTERIM INTEREST"]


def money(v):
    if v is None or v == "":
        return 0.0
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).strip().replace("$", "").replace(",", "").replace("(", "-").replace(")", "")
    try:
        return float(s) if s else 0.0
    except ValueError:
        return 0.0


def fmt_m(v: float, signed: bool = False) -> str:
    sign = ""
    if signed and v > 0:
        sign = "+"
    elif v < 0:
        sign = "-"
    return f"{sign}${abs(v)/1_000_000:.1f}M"


def norm_label(v) -> str:
    return re.sub(r"\s+", " ", str(v or "")).strip().upper()


def read_rows():
    wb = openpyxl.load_workbook(CAPEX_FILE, data_only=True)
    ws = wb["CapEx Spend by FY (2)"]
    return list(ws.iter_rows(values_only=True))


def parse_profile(rows, start_marker: str) -> dict[str, dict[str, float]]:
    start = None
    for i, row in enumerate(rows):
        if start_marker.lower() in str(row[1] or "").lower():
            start = i
            break
    if start is None:
        raise RuntimeError(f"Could not find section: {start_marker}")

    data: dict[str, dict[str, float]] = {}
    # Data begins two rows below marker after MODEL INPUTS/header row.
    # Stop at blank row or next section.
    for row in rows[start + 2:]:
        label = norm_label(row[1] if len(row) > 1 else "")
        if not label:
            break
        if label in {"TOTAL CAPEX BUDGET", "TOTAL EX INTERIM INTEREST", "LAND"}:
            continue
        if label not in ORDER:
            continue
        # Use the first occurrence of CONSTRUCTION/MHE as base spend; INTERIM INTEREST captures its own total.
        if label in data:
            continue
        data[label] = {fy: money(row[idx] if idx < len(row) else 0) for fy, idx in FY_COLS.items()}
    return data


def compute():
    rows = read_rows()
    original = parse_profile(rows, "Wallkill  - Original")
    updated = parse_profile(rows, "Wallkill Updated")
    cats = [c for c in ORDER if c in original or c in updated]
    out = []
    for c in cats:
        o = original.get(c, {fy: 0 for fy in FY_COLS})
        u = updated.get(c, {fy: 0 for fy in FY_COLS})
        delta = {fy: u.get(fy, 0) - o.get(fy, 0) for fy in FY_COLS}
        early_reduction = -sum(delta[fy] for fy in EARLY)  # positive means spend removed from FY27/FY28
        future_increase = sum(delta[fy] for fy in FUTURE)
        out.append({
            "category": c.title().replace("Mhe", "MHE"),
            "fy27": delta["FY27"],
            "fy28": delta["FY28"],
            "fy29": delta["FY29"],
            "fy30": delta["FY30"],
            "fy31": delta["FY31"],
            "early_reduction": early_reduction,
            "future_increase": future_increase,
            "interest_or_other": future_increase - early_reduction,
        })
    totals = {
        "fy27": sum(r["fy27"] for r in out),
        "fy28": sum(r["fy28"] for r in out),
        "fy29": sum(r["fy29"] for r in out),
        "fy30": sum(r["fy30"] for r in out),
        "fy31": sum(r["fy31"] for r in out),
        "early_reduction": sum(r["early_reduction"] for r in out),
        "future_increase": sum(r["future_increase"] for r in out),
    }
    totals["interest_or_other"] = totals["future_increase"] - totals["early_reduction"]
    return out, totals


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos Display"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    s.line.width = Pt(0.75)
    return s


def add_kpi(slide, x, title, value, subtitle, color):
    add_rect(slide, x, Inches(0.92), Inches(3.95), Inches(0.98), LIGHT_GRAY, MID_GRAY, True)
    add_text(slide, x + Inches(0.18), Inches(1.02), Inches(3.55), Inches(0.2), title, 9.5, True, GRAY)
    add_text(slide, x + Inches(0.18), Inches(1.25), Inches(1.7), Inches(0.35), value, 23, True, color)
    add_text(slide, x + Inches(1.85), Inches(1.29), Inches(1.85), Inches(0.26), subtitle, 9.5, False, BLACK)


def build():
    rows, totals = compute()

    prs = Presentation(str(TEMPLATE)) if TEMPLATE.exists() else Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Header
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.72), WMT_BLUE)
    add_text(slide, Inches(0.35), Inches(0.12), Inches(9.2), Inches(0.42),
             "Wallkill CapEx Pushout Supports Extended Permitting Timeline", 20, True, WHITE)
    add_text(slide, Inches(9.2), Inches(0.17), Inches(3.8), Inches(0.3),
             "FY27/FY28 cash draw rephased", 11, True, WMT_YELLOW, PP_ALIGN.RIGHT)

    add_kpi(slide, Inches(0.35), "FY27/FY28 spend reduced", fmt_m(totals["early_reduction"]),
            f"{fmt_m(-totals['fy27'], True)} FY27 / {fmt_m(-totals['fy28'], True)} FY28", RED)
    add_kpi(slide, Inches(4.68), "Moved into FY29-FY31", fmt_m(totals["future_increase"]),
            "aligned to later permit path", WMT_BLUE)
    add_kpi(slide, Inches(9.01), "Interim interest / total reduction", fmt_m(-totals["interest_or_other"]),
            "lower carrying cost vs. original", GREEN)

    # Waterfall-ish year delta chart
    add_text(slide, Inches(0.35), Inches(2.18), Inches(5.75), Inches(0.25),
             "Annual CapEx shift vs. original Wallkill profile", 12, True, BLACK)
    chart_data = CategoryChartData()
    chart_data.categories = ["FY27", "FY28", "FY29", "FY30", "FY31"]
    deltas = [totals["fy27"]/1_000_000, totals["fy28"]/1_000_000, totals["fy29"]/1_000_000, totals["fy30"]/1_000_000, totals["fy31"]/1_000_000]
    chart_data.add_series("$M change", deltas)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.35), Inches(2.5), Inches(5.75), Inches(3.75), chart_data).chart
    chart.has_legend = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.chart_title.has_text_frame = False
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = RED if deltas[i] < 0 else WMT_BLUE

    # Narrative box
    add_rect(slide, Inches(0.35), Inches(6.38), Inches(5.75), Inches(0.58), PALE_YELLOW, WMT_YELLOW, True)
    add_text(slide, Inches(0.52), Inches(6.47), Inches(5.4), Inches(0.34),
             "Message: CapEx moves out of FY27/FY28 and into FY29-FY31, reducing early cash deployment and interim interest while matching the extended permitting schedule.",
             9.5, True, BLACK)

    # Category table
    add_text(slide, Inches(6.45), Inches(2.18), Inches(6.45), Inches(0.25),
             "Category view: amount pushed from FY27/FY28 and receiving years", 12, True, BLACK)
    x0, y0 = Inches(6.45), Inches(2.5)
    col_w = [Inches(1.72), Inches(1.12), Inches(0.88), Inches(0.88), Inches(0.88), Inches(0.88)]
    headers = ["Category", "FY27/28 Out", "FY29", "FY30", "FY31", "Net Later"]
    row_h = Inches(0.34)
    x = x0
    for i, h in enumerate(headers):
        add_rect(slide, x, y0, col_w[i], row_h, WMT_DARK)
        add_text(slide, x + Inches(0.03), y0 + Inches(0.06), col_w[i] - Inches(0.06), Inches(0.18),
                 h, 7.6, True, WHITE, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
        x += col_w[i]

    y = y0 + row_h
    display_rows = sorted(rows, key=lambda r: abs(r["early_reduction"]), reverse=True)
    for idx, r in enumerate(display_rows):
        fill = WHITE if idx % 2 == 0 else RGBColor(247, 249, 252)
        vals = [r["category"], fmt_m(r["early_reduction"]), fmt_m(r["fy29"], True), fmt_m(r["fy30"], True), fmt_m(r["fy31"], True), fmt_m(r["future_increase"], True)]
        x = x0
        for i, val in enumerate(vals):
            add_rect(slide, x, y, col_w[i], row_h, fill, MID_GRAY)
            color = BLACK
            if i == 1:
                color = RED
            elif i >= 2:
                color = WMT_BLUE if not val.startswith("-") else RED
            add_text(slide, x + Inches(0.03), y + Inches(0.06), col_w[i] - Inches(0.06), Inches(0.18),
                     val, 7.4, i == 0, color, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
            x += col_w[i]
        y += row_h

    vals = ["TOTAL", fmt_m(totals["early_reduction"]), fmt_m(totals["fy29"], True), fmt_m(totals["fy30"], True), fmt_m(totals["fy31"], True), fmt_m(totals["future_increase"], True)]
    x = x0
    for i, val in enumerate(vals):
        add_rect(slide, x, y, col_w[i], row_h, PALE_BLUE, WMT_BLUE)
        add_text(slide, x + Inches(0.03), y + Inches(0.06), col_w[i] - Inches(0.06), Inches(0.18),
                 val, 7.8, True, WMT_DARK, PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT)
        x += col_w[i]

    add_text(slide, Inches(6.45), Inches(6.63), Inches(6.42), Inches(0.28),
             "Note: FY27/28 Out is calculated as the reduction from the original profile. Net Later is the increase/(decrease) across FY29-FY31. Interim Interest falls with the rephased cash draw.",
             7.1, False, GRAY)

    add_text(slide, Inches(0.35), Inches(7.12), Inches(12.5), Inches(0.18),
             "Source: Wallkill CapEx Push.xlsx, Original vs. Updated Wallkill CapEx profile.", 7.2, False, GRAY)

    prs.save(str(OUT))
    print(f"Created: {OUT}")
    print(f"FY27/28 reduction: {fmt_m(totals['early_reduction'])}")
    print(f"FY29-FY31 increase: {fmt_m(totals['future_increase'])}")
    print(f"Interest/total reduction: {fmt_m(-totals['interest_or_other'])}")


if __name__ == "__main__":
    build()
